"""
1단계: PPT 파일 전체 raw 추출 스크립트
- PPT 폴더 안의 모든 .pptx/.ppt 파일을 읽어서
- 텍스트 + 표 전부 추출하여
- raw_data 폴더에 파일별로 개별 JSON으로 저장
"""

import os
import json
from pptx import Presentation

PPT_FOLDER = "/Volumes/Untitled/USB Drive/티엔에프컨설팅/교육"
OUTPUT_FOLDER = "/Users/kyeongpilheo/Desktop/Python/proposal-ai/raw_data"

# 출력 폴더 생성
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

def extract_pptx(file_path):
    """PPT 파일에서 슬라이드별 텍스트 + 표 전부 추출"""
    prs = Presentation(file_path)
    slides = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            "슬라이드": slide_num,
            "텍스트": [],
            "표": []
        }

        for shape in slide.shapes:
            # 일반 텍스트 추출
            if hasattr(shape, "text") and shape.text.strip():
                if not shape.has_table:
                    slide_data["텍스트"].append(shape.text.strip())

            # 표 추출 (전체 셀 내용 보존)
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    # 빈 행은 제외
                    if any(cell for cell in row_data):
                        table_data.append(row_data)
                if table_data:
                    slide_data["표"].append(table_data)

        # 텍스트나 표 중 하나라도 있는 슬라이드만 저장
        if slide_data["텍스트"] or slide_data["표"]:
            slides.append(slide_data)

    return slides


def run():
    pptx_files = [
        f for f in os.listdir(PPT_FOLDER)
        if f.endswith(".pptx") or f.endswith(".ppt")
    ]

    if not pptx_files:
        print("❌ PPT 파일을 찾을 수 없습니다.")
        return

    print(f"📂 총 {len(pptx_files)}개 파일 발견\n")
    success, fail = 0, 0

    for i, filename in enumerate(pptx_files, 1):
        file_path = os.path.join(PPT_FOLDER, filename)
        # 저장 파일명: 번호_원본파일명.json
        output_filename = f"{i:03d}_{os.path.splitext(filename)[0]}.json"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)

        # 이미 처리된 파일은 스킵
        if os.path.exists(output_path):
            print(f"  ⏭️  [{i}/{len(pptx_files)}] 스킵 (이미 처리됨): {filename}")
            success += 1
            continue

        try:
            slides = extract_pptx(file_path)
            result = {
                "원본파일명": filename,
                "슬라이드수": len(slides),
                "슬라이드": slides
            }
            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(result, f, ensure_ascii=False, indent=2)

            print(f"  ✅ [{i}/{len(pptx_files)}] 완료: {filename} ({len(slides)}개 슬라이드)")
            success += 1

        except Exception as e:
            print(f"  ❌ [{i}/{len(pptx_files)}] 오류: {filename} → {e}")
            fail += 1

    print(f"\n{'='*50}")
    print(f"완료: {success}개 성공 / {fail}개 실패")
    print(f"저장 위치: {OUTPUT_FOLDER}")


if __name__ == "__main__":
    run()
